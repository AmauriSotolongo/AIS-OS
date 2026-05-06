#!/usr/bin/env python3
"""
Generador de propuestas Digital Compass — PPTX
Bike Method Phase 1: corre manualmente, revisa antes de enviar.
Adapted from The Three Ms of AI™ © 2026 Nate Herk.

Markup soportado en textos:
  **texto**  →  bold + italic + color plata (palabras clave destacadas)
"""

import argparse
import json
import os
import re
import sys
import tempfile
import urllib.request
from datetime import date
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Error: python-pptx no instalado. Corre: pip install python-pptx")
    sys.exit(1)

try:
    from openai import OpenAI
    OPENAI_AVAILABLE = True
except ImportError:
    OPENAI_AVAILABLE = False

# Load .env from project root
_env_path = Path(__file__).parent.parent / ".env"
if _env_path.exists():
    try:
        from dotenv import load_dotenv
        load_dotenv(_env_path)
    except ImportError:
        for line in _env_path.read_text().splitlines():
            if "=" in line and not line.startswith("#"):
                k, _, v = line.partition("=")
                os.environ.setdefault(k.strip(), v.strip())

# ── Color palette — Digital Compass (black & silver) ──────────────────────
DC_BLACK    = RGBColor(0x0D, 0x0D, 0x0D)   # fondo principal
DC_CHARCOAL = RGBColor(0x1A, 0x1A, 0x1A)   # fondo alterno
DC_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)   # texto principal
DC_SILVER   = RGBColor(0xC8, 0xC8, 0xC8)   # texto secundario / bullets
DC_PLATINUM = RGBColor(0xE8, 0xE8, 0xE8)   # highlights bold+italic
DC_MUTED    = RGBColor(0x77, 0x77, 0x77)   # fechas, notas

# Slide canvas — 1920×1080 widescreen
SLIDE_W = Emu(9144000)   # 10"
SLIDE_H = Emu(5143500)   # 5.625"


def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_design_accents(slide, variant: str = "default"):
    """Add silver geometric accents to slide background."""

    # ── Left accent bar — thick silver vertical strip
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(0.18), SLIDE_H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DC_SILVER
    bar.line.fill.background()

    # ── Bottom accent bar — full width
    bottom = slide.shapes.add_shape(
        1, Inches(0), Emu(SLIDE_H - Emu(55000)), SLIDE_W, Emu(55000)
    )
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = DC_MUTED
    bottom.line.fill.background()

    # ── Large background circle — outline only (decorative)
    cx, cy = (7.8, 2.8) if variant != "cover" else (7.5, 2.8)
    r = 2.2
    circ = slide.shapes.add_shape(
        9,  # MSO_SHAPE_TYPE oval
        Inches(cx - r), Inches(cy - r), Inches(r * 2), Inches(r * 2)
    )
    circ.fill.background()
    circ.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    circ.line.width = Emu(20000)

    # ── Small dot grid — top right corner
    for i in range(3):
        for j in range(3):
            dot = slide.shapes.add_shape(
                1,
                Inches(9.05 + j * 0.18), Inches(0.22 + i * 0.18),
                Inches(0.07), Inches(0.07)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = DC_SILVER if (i + j) % 2 == 0 else DC_MUTED
            dot.line.fill.background()



def _parse_markup(text: str):
    """
    Split text into segments: (content, is_highlight).
    **word** → is_highlight=True (bold + italic + platinum color).
    """
    parts = re.split(r'\*\*(.+?)\*\*', text)
    segments = []
    for i, part in enumerate(parts):
        if part:
            segments.append((part, i % 2 == 1))
    return segments


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=DC_WHITE,
                align=PP_ALIGN.LEFT, word_wrap=True):
    """Add a textbox with optional **markup** for bold+italic highlights."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align

    segments = _parse_markup(text)
    for content, is_highlight in segments:
        run = p.add_run()
        run.text = content
        run.font.size = Pt(font_size)
        run.font.bold = bold or is_highlight
        run.font.italic = is_highlight
        run.font.color.rgb = DC_PLATINUM if is_highlight else color
    return tb


def add_bullet_slide(slide, title, bullets, title_color=DC_WHITE, bullet_color=DC_SILVER):
    add_textbox(slide, title, 0.6, 0.9, 8.8, 0.65,
                font_size=26, bold=True, color=title_color)

    y = 1.65
    for bullet in bullets:
        add_textbox(slide, f"  {bullet}", 0.6, y, 8.8, 0.5,
                    font_size=17, color=bullet_color)
        y += 0.54


def generate_cover_image(proyecto: str, cliente: str) -> str | None:
    api_key = os.environ.get("OPENAI_API_KEY")
    if not api_key or not OPENAI_AVAILABLE:
        return None
    try:
        client = OpenAI(api_key=api_key)
        prompt = (
            f"Minimalist hero illustration for a professional services proposal. "
            f"Context: {proyecto} for {cliente}. "
            "Pure black background (#0D0D0D). Silver, platinum and white metallic tones only. "
            "Abstract geometric composition that visually represents the industry and service — "
            "use symbolic shapes, icons or textures related to the project theme. "
            "Clean lines, subtle silver glow, high contrast accents. "
            "No text, no logos, no people, no faces. "
            "High-end pitch deck cover aesthetic — cinematic, premium, dark mode. "
            "Monochromatic palette: black, charcoal, silver, platinum white."
        )
        response = client.images.generate(
            model="dall-e-3",
            prompt=prompt,
            size="1024x1024",
            quality="standard",
            n=1,
        )
        url = response.data[0].url
        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        urllib.request.urlretrieve(url, tmp.name)
        print(f"✓ Imagen IA generada")
        return tmp.name
    except Exception as e:
        print(f"⚠ Imagen no generada ({e}). Continuando sin imagen.")
        return None


def build_presentation(data: dict, cliente: str, proyecto: str, output_path: str):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank = prs.slide_layouts[6]
    fecha = date.today().strftime("%d de %B de %Y").replace(
        "January","enero").replace("February","febrero").replace(
        "March","marzo").replace("April","abril").replace(
        "May","mayo").replace("June","junio").replace(
        "July","julio").replace("August","agosto").replace(
        "September","septiembre").replace("October","octubre").replace(
        "November","noviembre").replace("December","diciembre")

    sc = data.get("slide_content", data)

    # ── SLIDE 1 — Portada ────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank)
    set_bg(s1, DC_BLACK)

    img_path = generate_cover_image(proyecto, cliente)
    text_w = 5.0 if img_path else 8.8

    if img_path:
        s1.shapes.add_picture(img_path, Inches(5.5), Inches(0), Inches(4.5), SLIDE_H)

    add_design_accents(s1, variant="cover")
    add_textbox(s1, "DIGITAL COMPASS", 0.6, 0.55, text_w, 0.4,
                font_size=11, bold=True, color=DC_MUTED)
    titulo = sc.get("portada", {}).get("titulo", "Propuesta de Proyecto")
    add_textbox(s1, titulo, 0.6, 1.3, text_w, 1.6, font_size=36, bold=True,
                color=DC_WHITE)
    add_textbox(s1, cliente, 0.6, 3.25, text_w, 0.55, font_size=20,
                color=DC_SILVER)
    add_textbox(s1, fecha, 0.6, 3.95, text_w, 0.4, font_size=13,
                color=DC_MUTED)

    # ── SLIDE 2 — Entendemos tu reto ─────────────────────────────────────
    s2 = prs.slides.add_slide(blank)
    set_bg(s2, DC_BLACK)
    add_design_accents(s2)
    add_textbox(s2, "Entendemos tu reto", 0.6, 0.9, 8.8, 0.65,
                font_size=28, bold=True)
    reto = sc.get("reto", {}).get("texto", "Descripción del reto.")
    add_textbox(s2, reto, 0.6, 1.9, 8.8, 3.0, font_size=19, color=DC_SILVER)

    # ── SLIDE 3 — Lo que proponemos ──────────────────────────────────────
    s3 = prs.slides.add_slide(blank)
    set_bg(s3, DC_CHARCOAL)
    add_design_accents(s3)
    add_bullet_slide(s3, "Lo que proponemos",
                     sc.get("propuesta", {}).get("bullets", []))

    # ── SLIDE 4 — Alcance ────────────────────────────────────────────────
    s4 = prs.slides.add_slide(blank)
    set_bg(s4, DC_BLACK)
    add_design_accents(s4)
    add_bullet_slide(s4, "Alcance del proyecto",
                     sc.get("alcance", {}).get("features", []))

    # ── SLIDE 5 — Tecnología ─────────────────────────────────────────────
    s5 = prs.slides.add_slide(blank)
    set_bg(s5, DC_CHARCOAL)
    add_design_accents(s5)
    add_bullet_slide(s5, "Tecnología y enfoque",
                     sc.get("tecnologia", {}).get("bullets", []))

    # ── SLIDE 6 — Timeline ───────────────────────────────────────────────
    s6 = prs.slides.add_slide(blank)
    set_bg(s6, DC_BLACK)
    add_design_accents(s6)
    add_textbox(s6, "Timeline estimado", 0.6, 0.9, 8.8, 0.65,
                font_size=28, bold=True)
    add_textbox(s6, "* Sujeto a revisión tras kick-off", 0.6, 1.65, 8.8, 0.38,
                font_size=12, color=DC_MUTED)
    y = 2.15
    for fase in sc.get("timeline", {}).get("fases", []):
        add_textbox(s6, f"  {fase}", 0.6, y, 8.8, 0.5,
                    font_size=17, color=DC_SILVER)
        y += 0.54

    # ── SLIDE 7 — Inversión ──────────────────────────────────────────────
    s7 = prs.slides.add_slide(blank)
    set_bg(s7, DC_CHARCOAL)
    add_design_accents(s7)
    inv = sc.get("inversion", {})
    add_textbox(s7, "Inversión", 0.6, 0.9, 8.8, 0.65, font_size=28, bold=True)
    add_textbox(s7, inv.get("total", "A cotizar"), 0.6, 2.0, 8.8, 0.9,
                font_size=52, bold=True, color=DC_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s7, inv.get("modalidad", "50% inicio · 50% entrega"),
                0.6, 3.2, 8.8, 0.55, font_size=17,
                color=DC_SILVER, align=PP_ALIGN.CENTER)

    # ── SLIDE 8 — Próximos pasos ─────────────────────────────────────────
    s8 = prs.slides.add_slide(blank)
    set_bg(s8, DC_BLACK)
    add_design_accents(s8)
    pasos = sc.get("proximos_pasos", {}).get("bullets", [])
    add_bullet_slide(s8, "Próximos pasos", pasos)
    contacto = sc.get("proximos_pasos", {}).get(
        "contacto", "digitalcompass.ia@gmail.com  |  +52 (998) 145 2628")
    add_textbox(s8, contacto, 0.6, 5.1, 8.8, 0.4, font_size=13,
                color=DC_MUTED, align=PP_ALIGN.CENTER)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    # Default output dir: ~/Desktop/Propuestas/
    prs.save(output_path)
    print(f"✓ Propuesta generada: {output_path}")

    if img_path:
        os.unlink(img_path)


def main():
    parser = argparse.ArgumentParser(description="Genera propuesta Digital Compass en PPTX")
    parser.add_argument("--cliente",  required=True)
    parser.add_argument("--proyecto", required=True)
    parser.add_argument("--output",   required=True)
    parser.add_argument("--data",     required=True)
    args = parser.parse_args()

    try:
        data = json.loads(args.data)
    except json.JSONDecodeError as e:
        print(f"Error en JSON: {e}")
        sys.exit(1)

    build_presentation(data, args.cliente, args.proyecto, args.output)


if __name__ == "__main__":
    main()
